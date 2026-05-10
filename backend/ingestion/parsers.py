"""Structured parsers for PrismAI ingestion.

Reads raw files and converts them into labelled structural elements.
This module never chunks — chunking is always a separate downstream step.
"""

# Future annotations means type hints are evaluated lazily as strings. This
# lets us write modern syntax like `str | None` and `list[tuple[...]]` even
# when the file is imported on older Python versions. Removing it would force
# us to import everything from `typing` and rewrite a few annotations.
from __future__ import annotations

# `re` powers the regex work in PDF heading detection and email cleanup
# (stripping HTML tags, recognising reply markers like "On ... wrote:"). If
# this import goes, both PDF parsing and email cleanup break.
import re

# `statistics.median` gives us the body-text font size for PDFs. We use the
# median (not the mean) so a single oversized title doesn't pull the baseline
# up and cause us to mis-classify normal body text as headings.
import statistics

# `dataclass` builds the ParsedElement record type for free; `field` lets us
# give the metadata dict a per-instance default. Without `field`, every
# ParsedElement would share the same dict object — a classic Python footgun
# where mutating one element's metadata would silently mutate every other.
from dataclasses import dataclass, field

# Standard typing imports. Kept small on purpose — modern Python lets us write
# `list[str]` directly, but Dict / List / Any improve readability for the PMs
# reading this file later.
from typing import Any, Dict, List

# PyMuPDF is imported as `fitz` because that's the historical name the library
# ships under. It is the fastest pure-Python PDF text extractor we've found
# and gives us per-span font sizes, which is how we tell headings from body
# text. Replacing it would mean rewriting parse_pdf entirely.
import fitz  # PyMuPDF

# `openpyxl` reads .xlsx files. We rely on its read-only / values-only mode in
# parse_excel so we don't load enormous workbooks fully into memory.
import openpyxl

# python-docx exposes Document at the top level, but to walk paragraphs and
# tables in their original document order we need the lower-level Paragraph
# and Table proxy classes too. Without them we'd lose table position relative
# to surrounding paragraphs, which downstream chunkers rely on.
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph

# python-pptx for PowerPoint. Only Presentation is needed at the entry point;
# slide / shape / notes objects are reached through it.
from pptx import Presentation


# A ParsedElement is the single unit of "structure" the rest of the pipeline
# understands. Every parser in this file produces a list of these (except the
# email and Excel parsers, which have specialised return shapes). Keeping the
# shape uniform means chunkers.py can treat docs, PDFs, slides, and text files
# with the same code paths.
@dataclass
class ParsedElement:
    """Carries one labelled piece of structure from a document before chunking.

    This is the shared contract between parsers.py and chunkers.py: every Word
    doc, PDF, slide deck, or plain-text/markdown file becomes an ordered list of
    ParsedElements so the rest of the pipeline never has to reopen the source file.

    Field roles (why each exists):
        - ``element_type``: Tells chunkers which routing rule applies — e.g.
          ``heading`` vs ``bold_title`` vs ``blank`` vs ``table`` vs ``slide``.
          Without a stable vocabulary, sprint-planning Word docs and exported PDF
          retrospectives could not share the same chunking logic.
        - ``content``: The extractable text for that unit (may be empty when the
          unit is purely structural, like a blank line). Chunkers join or split
          this text; embeddings store it.
        - ``level``: Heading depth (1–3) for hierarchy-aware chunking on OKR docs,
          PRDs, and meeting notes; ``0`` means "not a heading". We cap at 3 so Jira-
          export PDFs with ten nested headings do not explode metadata while rarely
          improving answer quality.
        - ``metadata``: Page numbers for PDFs, slide numbers for decks, row/column
          counts for Word tables — whatever helps citations and debugging. Uses
          ``field(default_factory=dict)`` so each element gets its own dict; a plain
          ``metadata={}`` on the class would make every row share one dict and corrupt
          citations across unrelated chunks.

    What breaks if this model is wrong: chunk_document would route to the wrong
    strategy, chunks would cite the wrong slide or page, or heading boundaries
    would vanish — so a PM asking "what did we decide in retro about OKR-3?" could
    get body text from the wrong section or no hit at all.
    """

    # `element_type` is the tag downstream chunkers route on. Restricting it
    # to a known vocabulary (the comment below documents the allowed values)
    # keeps detection logic in signal_detector.py and chunking in chunkers.py
    # from having to guess what a producer intended.
    element_type: str  # heading | paragraph | table | bold_title | blank | slide

    # The actual text payload. May be empty (blank rows, blank pages, etc.) —
    # downstream code must treat empty content as a structural marker, not as
    # data to embed.
    content: str

    # Heading level 1–3 (think H1/H2/H3 in Markdown or Word). Capped to 3
    # because anything deeper rarely changes retrieval quality but inflates
    # metadata. Non-headings always carry 0.
    level: int  # heading level 1–3, or 0 for non-headings

    # Free-form per-element metadata: page number for PDFs, slide number for
    # PPTX, row/col counts for tables. Using `field(default_factory=dict)` is
    # critical — a plain `= {}` would share one dict across every instance and
    # cross-contaminate metadata between unrelated elements.
    metadata: Dict[str, Any] = field(default_factory=dict)


# Word's XML uses namespaces, so a paragraph tag literally reads as
# "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}p". We just
# want the bare local name ("p", "tbl") to switch on. Without this helper we'd
# have to write that long namespace prefix everywhere we compare tags.
def _strip_xml_tag(tag: str) -> str:
    """Turn a namespaced Word XML tag into a short local name like ``p`` or ``tbl``.

    Word stores elements as ``{namespace}p`` and ``{namespace}tbl``; parse_docx must
    branch on plain ``p`` vs ``tbl`` so tables stay in reading order next to the
    paragraphs they illustrate (e.g. a capacity table under "Sprint 24 scope" in a
    planning doc).

    Decisions:
        1. Strip only the namespace suffix, not the whole string — regex on the full
           URI would be brittle across Office versions.
        2. Use ``rsplit("}", 1)[-1]`` not ``split("}")[-1]`` — if a pathological tag
           ever contained ``}`` twice, splitting from the right keeps the local name
           correct.
        3. If there is no ``}``, return the tag unchanged — handles odd fixtures and
           avoids turning valid local tags into empty strings.

    Returns the local tag string. parse_docx compares it to ``"p"`` and ``"tbl"``;
    a wrong return means no paragraphs or tables match and the PRD looks empty to
    retrieval. Missing this helper entirely would force duplicated namespace logic in
    every branch and guarantee drift bugs when Office updates namespaces.
    """
    # If a namespace is present, the local name is whatever comes after the
    # closing brace. `rsplit(..., 1)` is used so that a tag accidentally
    # containing `}` in the local part (extremely rare, but possible) doesn't
    # confuse us — we always split on the last `}` only.
    if "}" in tag:
        # RETURN: triggered when the XML element has a namespace (the normal
        # case for any Word .docx). Downstream, parse_docx compares this
        # against the literals "p" and "tbl" to route paragraphs vs tables.
        # If we returned the namespaced tag instead, NO branch in parse_docx
        # would ever match and the function would silently return an empty
        # list — every Word doc would look "empty" to the rest of the system.
        return tag.rsplit("}", 1)[-1]
    # RETURN: triggered when the tag has no namespace (e.g. test fixtures or
    # XML written by hand). Returning unchanged is correct because there is
    # nothing to strip; downstream callers receive the same tag they would
    # have gotten from a normal docx, so the routing logic still works.
    return tag


# Word style names look like "Heading 1", "Heading 2", "Title", etc. We want
# the integer level so chunkers can build a heading hierarchy. Doing this with
# a regex (instead of, say, a dict lookup) means we transparently handle
# slightly off-spec style names like "Heading2" or "Heading 12".
def _heading_level_from_style(style_name: str) -> int:
    """Extract a numeric heading level (1–3) from a Word paragraph style name.

    PM-authored PRDs and sprint summaries almost always use Word's built-in "Heading
    1/2/3" styles; chunkers need an integer so sections nest correctly under the right
    parent when someone asks "everything under Goals for Q3".

    Decisions:
        1. Regex-search for the first digit run instead of splitting on spaces — covers
           "Heading2", "Heading 12", and messy corporate templates without a fixed map.
        2. Clamp with ``max(1, min(3, n))`` — retrieval rarely benefits from H4–H9;
           deeper levels mostly add noise and break assumptions in signal_detector and
           chunkers that cap at three.
        3. Default to ``1`` when no digit appears (e.g. "Title") — treats cover titles
           like top-level anchors so the first chunk is not orphaned.

    Returns an int in 1..3. parse_docx puts this on ``ParsedElement.level`` for heading
    elements; wrong levels attach retro notes or OKR bullets to the wrong section in
    answers. Missing this function would force fragile string compares throughout.
    """
    # Search for the first run of digits anywhere in the style name.
    m = re.search(r"(\d+)", style_name)
    if m:
        # Clamp to 1–3 because deeper headings (H4+) rarely improve retrieval
        # but bloat metadata, and downstream chunkers assume max level 3.
        # `min(3, ...)` caps it from above; `max(1, ...)` floors it at 1
        # (defensive — the regex shouldn't yield 0, but if it did, level 0
        # would mean "this isn't a heading" and confuse callers).
        # RETURN: triggered when the style name contains a number (e.g.
        # "Heading 2"). Downstream, parse_docx writes this number into the
        # ParsedElement.level field, and chunkers use it to build a heading
        # hierarchy ("everything under H2 belongs to the parent H1"). A
        # wrong number here would put content under the wrong section in
        # retrieval and confuse the answers PMs see.
        return max(1, min(3, int(m.group(1))))
    # No digits found at all (e.g. style "Title") → treat as top-level.
    # RETURN: triggered when the style is "Title" / "Subtitle" / a custom
    # style name without a number. Defaulting to 1 means downstream chunking
    # treats these as top-of-document anchors, which is the right behaviour
    # for cover-page titles that PMs use as the doc's main heading.
    return 1


def parse_docx(file_path: str) -> List[ParsedElement]:
    """Convert a Word (.docx) file into an ordered list of structural elements.

    One-sentence summary: reads the document body in XML order and emits headings,
    paragraphs, bold titles, blanks, and tables as ParsedElements for chunking.

    Why it exists: PMs live in Word for PRDs, sprint planning one-pagers, stakeholder
    briefs, and retro write-ups. Those files mix narrative paragraphs with tables
    (scope matrices, RACI, acceptance criteria). The pipeline must preserve "this
    table belongs under this heading" or answers about Jira ticket tables end up
    detached from their owning section.

    Decisions inside (order matters):
        1. Walk ``doc.element.body`` instead of ``doc.paragraphs`` — the paragraphs
           API drops tables from the stream; ``doc.tables`` loses position relative to
           text. Body order is the only faithful representation for interleaved content.
        2. Wrap each ``p`` in ``Paragraph(element, doc._body)`` — python-docx requires
           the body parent to resolve styles and runs; skipping that breaks heading and
           bold detection.
        3. Heading branch on ``style_name.startswith("Heading")`` — Word's strongest
           structural signal; trusted over font alone for governance docs.
        4. Emit ``blank`` for whitespace-only paragraphs — chunkers use blanks as split
           hints for blank_line_based strategy; dropping them collapses sections.
        5. ``bold_title`` when all text-bearing runs are bold and length ≤ 80 — catches
           PMs who bold a single line instead of applying Heading styles; the 80-char
           cap avoids treating an entire bold disclaimer paragraph as a title.
        6. Keep ``paragraph`` content unstripped — preserves intentional internal spacing
           in acceptance criteria lists; headings/titles strip trailing Word noise.
        7. Tables: collapse consecutive duplicate cell text per row — Word repeats text
           for merged cells; without dedupe, embeddings see "Yes | Yes | Yes" instead of
           one "Yes".
        8. Join cells with `` | `` and rows with newlines — readable in chunks without
           introducing CSV commas that collide with comma-heavy prose.
        9. ``try/except`` → ``ValueError`` — corrupt or password-protected docs fail fast
           with a path the batch ingestor can log instead of crashing the whole job.

    Returns ``List[ParsedElement]``. Downstream: ``signal_detector.detect_docx_strategy``
    inspects structure; ``chunk_document`` groups by headings, blanks, or recursion. If
    this parser returned an empty list for a valid PRD, the PM would get no citations
    from that doc — silent gap in RAG for every Word-only artifact.

    Raises:
        ValueError: Wraps any underlying python-docx failure so callers can skip one bad file.
    """

    # Wrapping the whole parse in try/except keeps a single bad file from
    # taking down the ingestion pipeline. Any failure surfaces as a
    # ValueError with the file path, so the caller can log and skip cleanly.
    try:
        # Open the .docx. python-docx will raise immediately on a corrupt or
        # password-protected file, which we want — those should never reach
        # our element loop.
        doc = Document(file_path)
        # Output buffer. We accumulate parsed elements here and return at the
        # end so that on success the caller gets one consistent list.
        out: List[ParsedElement] = []

        # Walk the document *body* in raw XML order. This is the key trick:
        # iterating doc.paragraphs alone would silently drop tables, while
        # iterating doc.tables alone would lose their position relative to
        # surrounding paragraphs. Walking the body keeps everything
        # interleaved, which matters because PMs put tables in the middle of
        # PRDs — those tables need to chunk under the right heading.
        for element in doc.element.body:
            # Strip the XML namespace so "p" and "tbl" are easy to switch on.
            tag = _strip_xml_tag(element.tag)

            # Branch 1: this body element is a paragraph.
            if tag == "p":
                # Wrap the raw XML element in a Paragraph proxy so we get
                # convenient `.style`, `.text`, `.runs` access. Passing
                # `doc._body` as parent is required by python-docx's internal
                # API — without it the Paragraph can't resolve its style.
                para = Paragraph(element, doc._body)
                # Style may be None on stripped-down documents; default to ""
                # so downstream `startswith("Heading")` is safe.
                style_name = para.style.name if para.style is not None else ""
                # `.text` can be None on empty paragraphs; coerce to "" so we
                # never try to call string methods on None below.
                text = para.text or ""

                # Sub-branch 1a: paragraph is styled as a heading. Word
                # represents H1/H2/H3 with style names like "Heading 1". We
                # trust the style name as the strongest possible structural
                # signal — it was placed deliberately by the document author.
                if style_name.startswith("Heading"):
                    # Convert "Heading 2" → 2 (clamped to 1..3 in the helper).
                    lvl = _heading_level_from_style(style_name)
                    out.append(
                        ParsedElement(
                            element_type="heading",
                            # Strip whitespace because Word often pads
                            # heading text with trailing spaces from
                            # autocorrect.
                            content=text.strip(),
                            level=lvl,
                            metadata={},
                        )
                    )
                # Sub-branch 1b: paragraph contains nothing but whitespace.
                # We emit an explicit "blank" element rather than dropping it
                # because chunkers use blank lines as a signal to split
                # sections in some documents. Removing this branch would
                # collapse intentional blank lines and produce one giant
                # chunk per heading.
                elif not text.strip():
                    out.append(
                        ParsedElement(
                            element_type="blank",
                            content="",
                            level=0,
                            metadata={},
                        )
                    )
                # Sub-branch 1c: a styled-but-not-heading paragraph. We try
                # to recover faux-headings (PMs love to bold a single line
                # and call it a section title in Word).
                else:
                    # Build the list of *runs* (formatting segments) that
                    # actually contain text. A run with empty text usually
                    # means a tracked-changes artefact or a stray formatting
                    # toggle — including those would skew the all-bold check.
                    text_runs = [r for r in para.runs if (r.text or "")]
                    # "All bold" = there is at least one text-bearing run AND
                    # every text-bearing run is bold. `bool(text_runs)` is
                    # there to short-circuit `all([])` returning True (which
                    # would falsely flag empty paragraphs as bold titles).
                    all_bold = bool(text_runs) and all(r.bold for r in text_runs)
                    # Length cap of 80 chars filters out long bold passages
                    # (e.g. an entire paragraph styled bold for emphasis) so
                    # we only flag genuine title-like lines. Without this cap
                    # we'd treat highlighted body text as section breaks and
                    # over-fragment the document.
                    if all_bold and len(text) <= 80:
                        out.append(
                            ParsedElement(
                                element_type="bold_title",
                                content=text.strip(),
                                level=0,
                                metadata={},
                            )
                        )
                    # Default: ordinary paragraph. Note we keep the original
                    # `text` (no strip) so internal whitespace is preserved
                    # for chunkers that care about layout.
                    else:
                        out.append(
                            ParsedElement(
                                element_type="paragraph",
                                content=text,
                                level=0,
                                metadata={},
                            )
                        )

            # Branch 2: this body element is a table.
            elif tag == "tbl":
                # Wrap the XML element as a Table proxy for cell/row access.
                table = Table(element, doc._body)
                # We accumulate one display-string per row, then join.
                row_texts: List[str] = []
                # Track the widest row so we can record col_count in metadata.
                # Some Word tables have ragged rows when cells are merged or
                # added later, so we cannot rely on the first row's width.
                max_cols = 0
                for row in table.rows:
                    # Pull the row's cells once into a local variable.
                    # python-docx computes `row.cells` lazily on every
                    # access, so caching it once is both faster and
                    # protects us if the underlying XML mutates mid-loop
                    # (which can happen on documents with auto-fit tables).
                    cells = row.cells
                    # Update the running max column count; a simple `max()`
                    # is cheaper than computing it post-hoc from row_texts.
                    max_cols = max(max_cols, len(cells))
                    # `unique_vals` will hold the visible cell values in
                    # order, with consecutive duplicates removed.
                    unique_vals: List[str] = []
                    # `prev_val` tracks the previous cell's text so we can
                    # collapse "merged cell" duplicates. Word reports merged
                    # cells by repeating the same text in each spanned cell;
                    # without this dedupe the chunker would see "Yes | Yes |
                    # Yes" instead of the intended single "Yes".
                    prev_val: str | None = None
                    for cell in cells:
                        # Default to "" for cells with no text frame, then
                        # strip so trailing spaces don't break dedupe.
                        val = (cell.text or "").strip()
                        # Skip the cell if it matches the previous one — this
                        # is the merged-cell collapse described above.
                        if val == prev_val:
                            continue
                        unique_vals.append(val)
                        prev_val = val
                    # Format one row as "col1 | col2 | col3" — pipes are a
                    # safe separator because they almost never appear inside
                    # cell text and survive embedding well.
                    row_texts.append(" | ".join(unique_vals))

                # Stitch rows together with newlines so the table reads as
                # a small block of text downstream.
                table_body = "\n".join(row_texts)
                out.append(
                    ParsedElement(
                        element_type="table",
                        content=table_body,
                        level=0,
                        # row_count and col_count travel as metadata so the
                        # chunker can decide whether to keep a small table
                        # inline with its section or split it out.
                        metadata={
                            "row_count": len(table.rows),
                            "col_count": max_cols,
                        },
                    )
                )

        # RETURN: triggered when the entire body has been walked without
        # error. `out` is the ordered list of every heading, paragraph,
        # bold-title, blank, and table seen — exactly what chunkers.py
        # expects as input. The detection layer (signal_detector.py) and
        # the chunkers consume this list to produce the final embeddable
        # chunks. If this return ran with a half-built list (e.g. we broke
        # out of the loop early), downstream chunks would silently miss
        # whole sections of the original Word doc.
        return out

    # Catch *any* unexpected failure (corrupt zip container, missing parts,
    # XML errors) and re-raise as a clean ValueError. `from exc` preserves
    # the original traceback so debugging the underlying cause is still easy.
    except Exception as exc:
        # RAISE: triggered when python-docx blows up on a malformed file
        # (corrupt .docx, password-protected, ZIP-without-XML, etc.). The
        # ingestion pipeline catches this ValueError, logs the file path
        # we couldn't parse, and moves on to the next file — so a single
        # bad PRD never poisons a whole batch ingestion run.
        raise ValueError(f"Failed to parse DOCX {file_path!r}: {exc}") from exc


def parse_excel(file_path: str) -> Dict[str, Dict[str, Any]]:
    """Load every sheet of an Excel workbook into headers plus row dictionaries.

    One-sentence summary: streams each sheet, treats the first non-empty row as column
    headers, and returns one dict per sheet keyed by tab name for Excel-specific chunking.

    Why it exists: PMs track backlogs, OKR rollups, dependency matrices, and retro action
    items in .xlsx files (often exported from Jira or planning tools). Those grids are
    not prose; they need row-by-row or group-by-column chunking, so this parser outputs
    a tabular shape that ``chunk_excel_sheet`` consumes directly — not the same list type
    as Word/PDF.

    Decisions inside:
        1. ``read_only=True`` — workbooks can be huge; streaming avoids loading millions
           of cells into RAM during a batch ingest of quarterly OKR spreadsheets.
        2. ``data_only=True`` — embed calculated values (e.g. sum of story points) instead
           of raw ``=SUM(...)`` formula strings PMs cannot read in search results.
        3. ``try/finally`` with ``workbook.close()`` — read-only mode holds file handles;
           leaking them locks files on Windows and exhausts handles in long-running workers.
        4. Skip rows where every cell is None — Excel's "used range" often includes hundreds
           of blank rows at the bottom; keeping them would emit empty chunks.
        5. First non-empty row = headers — matches common PM sheets; we deliberately do not
           auto-detect title rows above headers (that needs product-specific rules).
        6. Synthetic ``Column{i}`` for blank header cells — guarantees every column has a
           stable key for ``row_dict`` and chunker ``f"{h}: {v}"`` lines.
        7. All cell values as strings — embeddings are text; mixing float/date types would
           complicate chunkers without improving recall on ticket exports.
        8. Skip data rows that are all whitespace after strip — removes visual spacer rows
           between epic blocks.
        9. Preserve sheet order via ``workbook.sheetnames`` — tab order often means
           "Summary" before "Details"; sorting alphabetically would scramble narrative flow.

    Returns ``{ sheet_name: {"headers": [...], "rows": [{col: str, ...}, ...]} }``.
    Downstream: pipeline calls ``chunk_excel_sheet`` per sheet with ``DetectionResult``;
    wrong shape here breaks grouping by epic or single-sheet OKR exports — PMs would see
    missing rows or mislabeled columns in answers.

    Raises:
        Nothing explicitly — openpyxl errors propagate unless callers wrap them; ``close()``
        still runs in ``finally``.
    """

    # `read_only=True` opens the workbook in streaming mode so we don't load
    # the full sheet into memory — important because Excel files PMs share
    # can be hundreds of MB. `data_only=True` returns the cached calculated
    # value of formulas (e.g. "=SUM(A1:A10)" → 42) rather than the formula
    # string. Without `data_only` we'd embed literal "=SUM(...)" text.
    workbook = openpyxl.load_workbook(
        file_path, read_only=True, data_only=True
    )
    # Wrap in try/finally so the workbook handle is always released, even on
    # an unexpected error halfway through. Leaving handles open in read-only
    # mode also keeps temp files locked on Windows.
    try:
        # Final return value: a dict keyed by sheet name. We build it up
        # incrementally so partial progress is preserved on any per-sheet
        # weirdness.
        result: Dict[str, Dict[str, Any]] = {}

        # Iterate sheets in their natural order in the workbook. PMs often
        # rely on tab order to imply meaning (overview → detail), so we
        # preserve it instead of, say, sorting alphabetically.
        for sheet_name in workbook.sheetnames:
            # Resolve the worksheet object for this tab name. Tab names can
            # collide with reserved words in theory, but openpyxl keys them
            # by string exactly as shown in Excel — we must use the same key
            # so "Sprint 24" rows never land under "Sheet1" by mistake.
            ws = workbook[sheet_name]
            # `raw_rows` collects every non-empty row for this sheet as a
            # tuple of cell values. Tuples (not lists) are used because the
            # streaming reader yields tuples and we don't need to mutate.
            raw_rows: List[tuple[Any, ...]] = []

            # Walk the sheet's rows. `values_only=True` returns plain Python
            # values (not Cell objects), keeping the loop fast and memory-
            # cheap on huge sheets.
            for row in ws.iter_rows(values_only=True):
                # Defensive: openpyxl can yield None for fully empty trailing
                # rows when the sheet has stale dimensions. Skip silently.
                if row is None:
                    continue
                # Skip rows where every cell is None — Excel files commonly
                # have hundreds of trailing blank rows from formatting that
                # would otherwise produce empty chunks downstream.
                if not any(cell is not None for cell in row):
                    continue
                raw_rows.append(tuple(row))

            # Sheet was completely empty (or only had whitespace rows). Emit
            # an empty record rather than skipping, so callers can tell the
            # difference between "sheet missing" and "sheet known empty".
            if not raw_rows:
                result[sheet_name] = {"headers": [], "rows": []}
                continue

            # First non-empty row is treated as the header row. We don't try
            # to detect "title" rows that PMs sometimes put above headers —
            # detecting that reliably needs domain knowledge that doesn't
            # belong in a generic parser. Document this expectation upstream.
            header_row = raw_rows[0]
            headers: List[str] = []
            for idx, cell in enumerate(header_row):
                # Synthesise a placeholder header when the cell is empty so
                # downstream code can always reference columns by *some*
                # name. Using the column index ("Column0", "Column1", ...)
                # keeps the placeholder unique within the sheet.
                if cell is None or str(cell).strip() == "":
                    headers.append(f"Column{idx}")
                else:
                    # Stringify even numeric headers — chunkers expect string
                    # keys when emitting "header: value" lines.
                    headers.append(str(cell))

            data_rows: List[Dict[str, str]] = []
            # Skip raw_rows[0] because that was the header row.
            for data_tuple in raw_rows[1:]:
                # Row dict is keyed by the header name, value is the string
                # form of the cell. We always store strings (never numbers,
                # dates, etc.) because everything downstream is going to be
                # embedded as text anyway.
                row_dict: Dict[str, str] = {}
                for col_idx, header in enumerate(headers):
                    # Some rows are shorter than the header row when the
                    # sheet has ragged trailing columns. Treat those missing
                    # cells as None / "".
                    val = (
                        data_tuple[col_idx]
                        if col_idx < len(data_tuple)
                        else None
                    )
                    # Empty cells become empty strings, never None — this
                    # keeps the chunker's "is value empty?" checks consistent
                    # without having to handle two different empty sentinels.
                    row_dict[header] = "" if val is None else str(val)

                # Drop rows whose *every* value is empty after stripping.
                # These are typically spacer rows added by the author for
                # visual grouping; keeping them would create empty chunks.
                if not any(v.strip() for v in row_dict.values()):
                    continue
                data_rows.append(row_dict)

            # Store the per-sheet result. The shape mirrors what
            # chunkers.chunk_excel_sheet expects, so no transformation is
            # needed downstream.
            result[sheet_name] = {"headers": headers, "rows": data_rows}

        # RETURN: triggered after every sheet has been processed without
        # error. The dict carries one entry per tab in the original
        # workbook, with header/row data ready for chunkers.chunk_excel_sheet
        # to consume. Excel does NOT go through chunk_document — pipeline.py
        # iterates this dict and calls the Excel-specific chunker per sheet,
        # so an empty `result` here would silently produce zero chunks for a
        # supposedly populated workbook.
        return result

    # `finally` guarantees `close()` runs even if an exception propagates.
    # openpyxl in read-only mode leaves a streaming reader holding the file;
    # not closing causes file-handle leaks on long-running processes.
    finally:
        workbook.close()


def parse_pptx(file_path: str) -> List[ParsedElement]:
    """Turn each slide in a PowerPoint deck into one ``slide`` ParsedElement in order.

    One-sentence summary: for every slide, merges title, body text shapes, and speaker
    notes into a single string and tags it with a 1-based slide number.

    Why it exists: Roadmaps, exec readouts, and sprint review decks are how PMs communicate
    narrative when Jira tickets are not enough. Retrieval must answer "what did we say on
    slide 7 about dependencies?" — so one element per slide matches how humans reference decks.

    Decisions inside:
        1. Title first, then other shapes with text frames — reading order matches how an
           audience scans the slide (title → bullets).
        2. Skip shapes without ``has_text_frame`` — charts and images have no extractable
           text without OCR; silently skipping avoids crashes; image-heavy decks need a
           separate OCR path documented upstream.
        3. Skip ``shape is title_shape`` in the body loop — prevents duplicating the title
           in the chunk and wasting context window.
        4. Strip non-empty paragraphs only — removes empty bullet placeholders left from
           corporate templates.
        5. Append speaker notes when present — notes often hold the real story (talk track)
           while the slide is minimal; omitting notes loses half the PM intent on review decks.
        6. ``slide_number = slide_idx + 1`` — humans say "slide 12", not "slide 11 in zero-
           based indexing"; metadata feeds citation UI and chunk metadata.
        7. Always emit one ParsedElement per slide even if ``content`` is empty — keeps slide
           indices aligned with the file; empty slides still occupy a slot in the deck order.

    Returns ``List[ParsedElement]`` with ``element_type="slide"``. Downstream:
    ``Strategy.slide_based`` + ``chunk_slide_elements`` produce exactly one chunk per slide.
    If this returned fewer elements than slides (e.g. skipping blank slides), citations would
    point at wrong slide numbers and PMs would distrust answers for quarterly business reviews.

    Raises:
        Propagates python-pptx errors (corrupt pptx) — callers should catch like other parsers.
    """

    # Open the deck. python-pptx keeps the file open lazily, but for our
    # one-pass read we don't need to manage the handle explicitly.
    prs = Presentation(file_path)
    elements: List[ParsedElement] = []

    # Iterate slides in their natural order. `enumerate` gives us a 0-based
    # index that we convert to 1-based for human-readable slide numbers.
    for slide_idx, slide in enumerate(prs.slides):
        # `parts` collects title, body, and notes for this slide. We join
        # them at the end so the slide reads top-to-bottom in the chunk.
        parts: List[str] = []

        # Pull the title shape first. Slides may legitimately have no title
        # (e.g. a section divider), so this can be None.
        title_shape = slide.shapes.title
        # Only proceed if there is a title AND it actually holds text. Some
        # title placeholders are picture- or chart-based and would crash on
        # `.text_frame` access without the has_text_frame check.
        if title_shape is not None and title_shape.has_text_frame:
            t = title_shape.text_frame.text.strip()
            # Skip blank titles — appending an empty string would create a
            # leading blank line in the chunk and waste embedding budget.
            if t:
                parts.append(t)

        # Walk every shape on the slide for body content.
        for shape in slide.shapes:
            # Pictures, charts, smart-art, etc. don't have a text frame — we
            # can't extract their text without OCR, so skip silently. (PM
            # slides are heavy on these; documenting upstream that decks
            # with image-only slides need a separate OCR step.)
            if not shape.has_text_frame:
                continue
            # Don't re-emit the title — we already added it as the first
            # part. `is` (identity) is correct here because python-pptx
            # returns the same shape object both ways.
            if title_shape is not None and shape is title_shape:
                continue
            # Inside a text frame, paragraphs map to bullet lines. We strip
            # each paragraph and drop empty ones to compress runs of blank
            # bullets that PMs leave behind from template editing.
            paras = [
                p.text.strip()
                for p in shape.text_frame.paragraphs
                if p.text and p.text.strip()
            ]
            # Join bullets with single newlines so they read as a list, then
            # add as one part. Skipping the append entirely on empty `paras`
            # keeps the slide chunk tight.
            if paras:
                parts.append("\n".join(paras))

        # Speaker notes are separately structured. They're huge for PM
        # decks because notes often contain the actual narrative the deck
        # only sketches.
        if slide.has_notes_slide:
            ns = slide.notes_slide
            # `notes_text_frame` can be None for empty notes panes.
            if ns.notes_text_frame is not None:
                ntxt = ns.notes_text_frame.text.strip()
                # Append notes only if there's something to append. Same
                # reasoning as title/body — no empty-string parts.
                if ntxt:
                    parts.append(ntxt)

        # Stitch the parts together with newlines so title/body/notes are
        # visually separated in the resulting chunk content.
        content = "\n".join(parts)
        elements.append(
            ParsedElement(
                element_type="slide",
                content=content,
                level=0,
                # 1-indexed slide number is what humans cite in reviews
                # ("see slide 12"). 0-indexed would confuse downstream UI
                # without a re-translation step.
                metadata={"slide_number": slide_idx + 1},
            )
        )

    # RETURN: triggered after every slide has been processed. `elements` is
    # one ParsedElement per slide, in deck order, with title/body/notes
    # already merged into a single `content` string. Downstream, chunkers
    # will route this list via Strategy.slide_based, which produces exactly
    # one chunk per slide. If we returned an empty list here on a populated
    # deck, the deck would be invisible to retrieval — PMs searching for
    # "Q3 launch slide" would get nothing.
    return elements


# PDFs don't expose a clean "is bold" property. We have to infer it from
# either the font name (e.g. "Helvetica-Bold") or a bitmask that PyMuPDF
# packs into the `flags` integer. Both signals are checked because a font
# may be named generically ("ArialMT") while still rendering bold via the
# flag, and vice versa.
def _span_is_bold(span: Dict[str, Any]) -> bool:
    """Infer whether a PDF text span was rendered bold using PyMuPDF metadata.

    One-sentence summary: returns True if the font name advertises bold or if PyMuPDF's
    span ``flags`` bitmask has the bold bit set.

    Why it exists: exported PDFs from Google Docs, Confluence, or emailed PRDs rarely expose
    Word-style heading tags — bold plus font size is often the only cue that "Risk register"
    is a subheading versus body copy in a dense retrofit summary.

    Decisions inside:
        1. Check font substring ``"Bold"`` / ``"bold"`` — many generators encode bold in the
           PostScript name (``Helvetica-Bold``). Cheap and works across locales that reuse Latin names.
        2. Fall back to ``flags & 16`` — bit 4 is MuPDF's bold marker when the font name stayed
           generic (``ArialMT``). Using ``isinstance(flags, int)`` avoids crashing on malformed PDFs.
        3. Both checks are needed — half of enterprise PDFs lie on one signal or the other;
           relying on only one collapses heading detection for retro packets.

    Returns ``bool``. Used only by ``parse_pdf`` for the "bold short line near body size"
    heading heuristic. Returning True everywhere would flood chunks with false headings;
    returning False everywhere would flatten structure — PMs asking scoped questions would get
    wall-of-text answers with no section boundaries.

    Raises:
        Never raises — defensive ``get`` / ``isinstance`` keep PDF oddities from bubbling up here.
    """

    # Default to "" so substring checks are always safe.
    font = span.get("font") or ""
    # Substring check covers both common spellings without normalising the
    # whole string (cheap, and avoids slicing edge cases). Removing this
    # would miss ~half of real-world PDFs because PyMuPDF's flag field is
    # often unreliable on poorly tagged documents.
    if isinstance(font, str) and ("Bold" in font or "bold" in font):
        # RETURN: triggered when the font name itself advertises bold
        # (e.g. "Helvetica-Bold", "ArialMT-BoldMT"). The caller (parse_pdf)
        # uses this True as part of its heading rule — a bold short span
        # near body size becomes a subheading. Without this signal, dense
        # PRDs that use bold subheadings instead of bigger fonts would
        # have their entire structure invisible to chunking.
        return True
    flags = span.get("flags")
    # Bit 4 (decimal 16) is PyMuPDF's "bold" indicator. The `& 16` mask
    # isolates that single bit so we don't false-positive on other flags
    # like italic or serif.
    if isinstance(flags, int) and (flags & 16):
        # RETURN: triggered when PyMuPDF's flags integer has the bold bit
        # set even though the font name didn't advertise it. This is the
        # backup signal — sometimes PDFs use a generic font like "ArialMT"
        # but render bold via the flag. parse_pdf uses this exactly the
        # same way as the font-name signal above.
        return True
    # RETURN: triggered when neither the font name nor the flags indicate
    # bold. parse_pdf treats this span as plain body text; if we were
    # wrong here on a real heading, that heading would be misclassified
    # as a paragraph and the chunker would lose a section break.
    return False


def parse_pdf(file_path: str) -> List[ParsedElement]:
    """Extract headings and paragraphs from a PDF using font-size and bold heuristics.

    One-sentence summary: estimates body font size from the first pages, then walks every page's
    text spans and labels each as heading or paragraph (with OCR placeholders for scanned pages).

    Why it exists: Legal PDFs, exported Confluence pages, emailed specs, and board decks are often
    PDF-only. PMs still need to query them like Word docs — "pull the acceptance criteria from the
    PDF attachment" — without manual copy-paste.

    Decisions inside:
        1. Two-pass design — pass 1 collects font sizes from first ``min(5, page_count)`` pages only
           for speed on 200-page contracts while the template font still appears early.
        2. ``get_text("dict")`` not plain ``get_text()`` — only dict mode exposes per-span ``size``
           and ``flags`` needed to separate H1-sized titles from body text.
        3. Ignore non-text blocks (``type != 0``) — image blocks lack ``lines``; skipping avoids noise.
        4. ``statistics.median`` for ``body_size`` not mean — one 28pt title must not pull the baseline
           up and mark every 11pt paragraph as a "small heading".
        5. Default ``body_size = 11.0`` when no sizes found — typical body point size; scanned PDFs then
           mostly fall through to paragraph or the short-page OCR branch.
        6. Short-page rule ``len(page_text.strip()) < 50`` — likely scanned or image-only; emit a single
           placeholder paragraph so ingestion does not silently drop pages (PM would otherwise think the
           PDF was indexed when it was empty).
        7. Heading if ``size >= 1.3 * body`` OR (bold AND ``len < 100`` AND size within 10% of body) —
           catches oversized titles and "bold subhead" patterns common in PRDs without changing font size.
        8. Level 1 vs 2 from ``1.5 * body`` threshold — coarse but stable; then ``max(1, min(3, lvl))``
           matches the rest of the pipeline's three-level cap for OKR-style outlines.
        9. Trailing ``blank`` with ``page_break_after`` — gives chunkers a soft boundary when headings
           are missing (meeting minutes PDFs).
        10. ``with fitz.open`` — ensures file handles release on Windows so temp uploads are not locked.

    Returns ``List[ParsedElement]`` in reading order with ``metadata["page"]`` on text elements.
    Downstream: ``chunk_document`` + detection may use heading_based or recursive_fallback; wrong output
    yields wrong citations ("page 14" vs actual) or no structure — PM gets irrelevant excerpts from legal
    PDFs or misses scanned appendix pages entirely without the OCR hint.

    Raises:
        Propagates PyMuPDF open errors for corrupt paths — callers should wrap batch jobs.
    """

    # Accumulator for every emitted span, OCR placeholder, and page-break marker
    # across the full PDF. Built incrementally page-by-page — we cannot pre-size it
    # because scanned vs text pages produce different element counts; truncating early
    # would drop appendix pages from board packs or compliance PDFs.
    elements: List[ParsedElement] = []

    # `with fitz.open(...)` guarantees the PDF is closed even if parsing
    # raises mid-loop. PDF readers hold native handles and need to be
    # released explicitly on Windows or the file stays locked.
    with fitz.open(file_path) as pdf:
        # ── First pass: figure out what "body text size" means for this PDF.
        # Different PDFs use different baselines (10pt for legal docs, 12pt
        # for slides exported as PDF), so we can't hardcode a number.
        font_sizes_first_pass: List[float] = []
        # Sample only the first few pages for speed. Large reports can have
        # hundreds of pages, but the body font is set in the template and
        # almost always visible by page 5.
        page_limit = min(5, pdf.page_count)
        for pi in range(page_limit):
            page = pdf[pi]
            # `get_text("dict")` returns the structured layout: blocks → lines
            # → spans, each span carrying its own font size and flags. This
            # is the only mode that exposes per-span font sizes — plain
            # `get_text()` returns just text strings.
            page_dict = page.get_text("dict")
            for block in page_dict.get("blocks", []):
                # block type 0 = text; type 1 = image. We only want text.
                # Without this guard, image blocks would have no `lines`
                # field and the inner loop would silently iterate nothing,
                # but explicit is better than implicit here.
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        sz = span.get("size")
                        # Defensive isinstance — PyMuPDF *should* always
                        # return a float, but a malformed PDF can yield None
                        # or a string and we'd rather skip than crash.
                        if isinstance(sz, (int, float)):
                            font_sizes_first_pass.append(float(sz))

        # Use the median, not the mean. A single H1 at 24pt would pull a
        # mean baseline up enough to mis-classify all 11pt body text as
        # body+, breaking heading detection. The median is robust to those
        # outliers because half the spans are above and half are below.
        if font_sizes_first_pass:
            body_size = float(statistics.median(font_sizes_first_pass))
        else:
            # No font info at all (unusual; usually means a scanned PDF).
            # 11pt is the most common body size in PM documents and a
            # reasonable default. Heading detection later degrades
            # gracefully if our guess is wrong because we'll just fall
            # through to "everything is a paragraph".
            body_size = 11.0

        # ── Second pass: emit elements page by page across the full deck.
        num_pages = pdf.page_count
        for pi in range(num_pages):
            page = pdf[pi]
            # 1-based page numbers for human-friendly metadata ("p. 14").
            page_num = pi + 1
            # Plain text length, used only as an OCR detector. We strip to
            # ignore whitespace-only pages.
            page_text = page.get_text()
            text_len = len(page_text.strip())

            # If the page has almost no extractable text (<50 chars), it's
            # almost certainly a scanned image. We emit a single placeholder
            # paragraph flagging that OCR is needed, and an explicit blank
            # marker for the page break, then move on. Without this branch
            # the parser would emit nothing for scanned pages and the user
            # would silently lose entire pages of content.
            if text_len < 50:
                elements.append(
                    ParsedElement(
                        element_type="paragraph",
                        content=f"[Page {page_num} — scanned, needs OCR]",
                        level=0,
                        metadata={"page": page_num},
                    )
                )
                # The trailing blank carries `page_break_after` so the
                # chunker can use it as a natural section break.
                elements.append(
                    ParsedElement(
                        element_type="blank",
                        content="",
                        level=0,
                        metadata={"page_break_after": page_num},
                    )
                )
                continue

            # Real text page → walk the structured layout and decide per
            # span whether it's a heading or a body line.
            page_dict = page.get_text("dict")
            for block in page_dict.get("blocks", []):
                # Skip image blocks (same reason as in the first pass).
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text = span.get("text") or ""
                        # Skip whitespace-only spans early — emitting them
                        # would create useless empty paragraphs.
                        if not text.strip():
                            continue
                        sz_raw = span.get("size")
                        # Same defensive isinstance as above. If we don't
                        # have a usable size we can't classify, so skip.
                        if not isinstance(sz_raw, (int, float)):
                            continue
                        sz = float(sz_raw)
                        bold = _span_is_bold(span)
                        # Heading rule:
                        #   • clearly bigger than body (≥ 1.3× median)
                        #     OR
                        #   • bold AND short (<100 chars) AND about the
                        #     same size as body (within 10%) — captures
                        #     "bold subheadings" used in dense PM PRDs
                        #     where authors don't bump the font size.
                        # Both conditions are needed: relying only on size
                        # misses bold-only subheads; relying only on bold
                        # over-flags emphasised body text.
                        is_heading = sz >= 1.3 * body_size or (
                            bold
                            and len(text) < 100
                            and abs(sz - body_size) <= 0.10 * body_size
                        )

                        if is_heading:
                            # H1 if dramatically larger; otherwise H2. This
                            # is intentionally coarse — clamping into 1..3
                            # below means deeper hierarchies get flattened.
                            lvl = 1 if sz >= 1.5 * body_size else 2
                            elements.append(
                                ParsedElement(
                                    element_type="heading",
                                    content=text,
                                    # Clamp 1..3 to match the rest of the
                                    # pipeline. min/max combo is the
                                    # idiomatic Python clamp.
                                    level=max(1, min(3, lvl)),
                                    metadata={"page": page_num},
                                )
                            )
                        else:
                            elements.append(
                                ParsedElement(
                                    element_type="paragraph",
                                    content=text,
                                    level=0,
                                    metadata={"page": page_num},
                                )
                            )

            # End-of-page marker. Always emitted (even after a real text
            # page) so chunkers can use page boundaries as a soft hint
            # when no headings are present.
            elements.append(
                ParsedElement(
                    element_type="blank",
                    content="",
                    level=0,
                    metadata={"page_break_after": page_num},
                )
            )

    # RETURN: triggered after every page has been processed (including
    # OCR-needed scanned pages, which were emitted as placeholders). The
    # caller routes this list through chunk_document; if there were
    # detected headings it goes via heading_based, otherwise via the
    # recursive fallback. An empty list here would mean the entire PDF
    # is invisible to retrieval — every search across that PDF would
    # come back empty even though the file was successfully ingested.
    return elements


def parse_text_file(file_path: str) -> List[ParsedElement]:
    """Parse a UTF-8 text or Markdown file line-by-line into ParsedElements.

    One-sentence summary: reads all lines, classifies each as blank, markdown heading, or paragraph.

    Why it exists: Roadmaps, OKRs, ADRs, and retro notes often live in ``.md`` or ``.txt`` in repos
    before they become Word docs. Line-oriented parsing mirrors how ``signal_detector.detect_text_strategy``
    counts ``#`` headings and blank-line ratios — parser output must align with those signals.

    Decisions inside:
        1. ``encoding="utf-8", errors="ignore"`` — PM tooling (Notion export, Slack paste) injects
           invalid bytes; ignoring preserves the rest of the file instead of failing the whole ingest.
        2. ``strip()`` per line — removes trailing ``\\n`` and leading noise; both classification and
           stored content use the stripped form so chunk boundaries match detector logic.
        3. Empty stripped line → ``blank`` — enables blank_line_based chunking for notes separated by
           empty lines (common in meeting minutes).
        4. Leading ``#`` run → heading — count consecutive ``#`` only at the start (not ``count("#")``)
           so inline hashes in titles like ``# OKR-3 #launch`` do not inflate level.
        5. ``level = max(1, min(3, level))`` — same three-level cap as Word/PDF so chunkers do not see
           unsupported H5 depths from over-eager markdown.
        6. Non-heading lines → ``paragraph`` — one physical line per element; soft wrapping in source
           becomes separate elements (acceptable trade-off for simple files).

    Returns ``List[ParsedElement]``. Downstream: ``chunk_document`` after ``detect_text_strategy`` chooses
    heading_based, blank_line_based, or recursive_fallback. Wrong classification would mis-route OKR
    markdown — e.g. every line treated as heading — producing tiny useless chunks and nonsense answers.

    Raises:
        File I/O errors propagate (missing file, permissions) — batch pipeline should catch per path.
    """

    # `errors="ignore"` silently drops bytes that aren't valid UTF-8. PMs
    # paste from Word / Notion / Slack all the time, which sometimes leaves
    # smart quotes or invisible control characters in supposedly-plain
    # files. We'd rather lose a few odd bytes than crash on a whole file.
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        lines = f.readlines()

    # Running output for this file — order matches top-to-bottom reading order in the
    # repo file so sprint-review notes above "### Risks" stay before that heading in the list.
    # chunk_document relies on sequence; reversing or sorting would scramble OKR markdown chunks.
    out: List[ParsedElement] = []
    for raw_line in lines:
        # Strip leading/trailing whitespace once. We use the stripped form
        # for both detection AND content because trailing newlines and
        # leading indentation rarely carry meaning in chunked text.
        line = raw_line.strip()

        # Empty lines become blank markers. The chunker uses runs of blanks
        # as candidate split points.
        if not line:
            out.append(
                ParsedElement(element_type="blank", content="", level=0)
            )
            continue

        # Markdown-style headings start with one or more `#` characters.
        if line.startswith("#"):
            # Counter for the leading `#` symbols. We initialise to 0 and
            # increment as we walk the prefix — starting elsewhere (1, say)
            # would give a wrong heading level, which would put content
            # under the wrong section in retrieval.
            level = 0
            # Walk characters until we hit a non-`#`. We can't just use
            # `line.count("#")` because that would count `#` characters
            # that appear *after* the heading text too (e.g. "# Section #1").
            for ch in line:
                if ch == "#":
                    level += 1
                else:
                    break
            # Clamp 1..3 to match the rest of the pipeline. A `# # # # #`
            # H5 would otherwise produce level=5 and confuse the chunker.
            level = max(1, min(3, level))
            # Strip the leading `#`s and the space that conventionally
            # follows them. We slice by `level` (the count of hashes) so
            # this works even if there's no space after the hashes.
            inner = line[level:].strip()
            out.append(
                ParsedElement(
                    element_type="heading",
                    content=inner,
                    level=level,
                )
            )
        else:
            # Anything that's not blank and not a heading is body text.
            # Keep the stripped form to suppress leading indentation that
            # would otherwise appear inside chunks for no semantic reason.
            out.append(
                ParsedElement(
                    element_type="paragraph",
                    content=line,
                    level=0,
                )
            )

    # RETURN: triggered after every line in the file has been classified
    # into blank / heading / paragraph. The list goes to chunk_document,
    # which dispatches via heading_based (if signal_detector saw enough
    # markdown headings) or blank-line / recursive fallback. An empty list
    # here would mean a perfectly readable .md file gets ingested as zero
    # chunks — silent data loss for any roadmap or OKR doc that lives in
    # markdown.
    return out


def parse_email_body(raw_body: str, is_html: bool = False) -> str:
    """Normalize an email body to plain text by stripping markup and quoted threads.

    One-sentence summary: optionally removes HTML cruft, then cuts everything from the first classic
    reply marker onward, then collapses excessive whitespace.

    Why it exists: PM workflows run on Gmail threads — sprint decisions, Jira ticket discussions, exec
    forwards. Raw bodies contain HTML wrappers, CSS, and full quoted history. Embedding that verbatim
    duplicates the same decision across chunks (once per reply), pollutes similarity search with footer
    boilerplate, and wastes tokens on ``<div>`` noise.

    Decisions inside:
        1. Normalize ``\\r\\n`` → ``\\n`` first — Outlook vs Gmail line endings otherwise break regexes
           anchored to line starts (false negatives on "On … wrote:" detection).
        2. HTML branch gated on ``is_html`` — plain-text bodies may contain ``<`` (code snippets, "a<b");
           stripping tags unconditionally would delete legitimate characters.
        3. Remove ``<style>`` and ``<script>`` before generic tag strip — keeps embedded CSS/JS from becoming
           gibberish sentences in chunks.
        4. ``re.sub(r"<[^>]+>", "", text)`` after — removes remaining tags; done after blocks so inner angle
           brackets in scripts do not leak.
        5. ``&#\\d+;`` → empty — avoids half-decoded numeric entities; trade-off: may drop a curly quote but
           prevents broken unicode from dominating embeddings.
        6. Named entities ``&nbsp;`` ``&amp;`` ``&lt;`` ``&gt;`` — restores readable spacing and symbols for
           retro summaries pasted from webmail.
        7. Reply cut uses line-scanned patterns (On … wrote:, Original Message, ``>``, underscores,
           From:/Sent: pair) — first match wins so only the newest composer message survives; ``From:`` alone
           does not cut (avoids killing lines like "From: the team — thanks!").
        8. Collapse ``\\n{3,}`` and runs of spaces — signatures add vertical whitespace that would become
           empty or near-empty chunks downstream.

    Returns a single ``str`` (not ParsedElements). Downstream: ``chunk_email`` prefixes From/Subject/Date and
    splits on paragraphs. If this returned raw HTML or full threads, a PM asking "what did Sarah decide about
    the milestone?" could retrieve six copies of the same thread or random footer legal text instead of the
    latest reply.

    Raises:
        Never raises — regex replace on bad input yields best-effort text.
    """

    # Normalise Windows-style line endings to Unix newlines first. Outlook
    # emits \r\n, Gmail emits \n, and mixing them confuses the regex line
    # checks below (especially `\s*$` anchors). Doing this once up front
    # means every later step sees consistent line breaks.
    text = raw_body.replace("\r\n", "\n")

    # ── HTML cleanup. Only run if the caller flagged the body as HTML —
    # forcing this on plain-text emails would chew through perfectly fine
    # text that happens to contain `<` characters (e.g. code snippets,
    # "<see attached>").
    if is_html:
        # Remove <style> blocks *including their contents*. Without this,
        # CSS rule text would survive and pollute the chunk with junk like
        # ".x_section{font-family:Arial;...}". `re.DOTALL` makes `.` match
        # newlines (style blocks span lines); IGNORECASE handles tag-case
        # variations from different mail clients.
        text = re.sub(
            r"<style[^>]*>.*?</style>",
            "",
            text,
            flags=re.DOTALL | re.IGNORECASE,
        )
        # Same treatment for <script>. Less common in email than style
        # blocks but still occurs in marketing emails and bug reports.
        text = re.sub(
            r"<script[^>]*>.*?</script>",
            "",
            text,
            flags=re.DOTALL | re.IGNORECASE,
        )
        # Strip every remaining HTML tag. Done *after* style/script removal
        # so we don't leave their textual content behind. The greedy `[^>]+`
        # handles attributes; switching to a lazy match would break on tags
        # containing `>` inside attribute values (rare but real).
        text = re.sub(r"<[^>]+>", "", text)
        # Drop numeric character references like `&#8217;` outright. We
        # don't try to translate them to characters because mail clients
        # use them inconsistently and the cost of mistranslation outweighs
        # the cost of dropping a curly quote.
        text = re.sub(r"&#\d+;", "", text)
        # Decode the four named entities that account for ~99% of real
        # email HTML. Order matters: `&amp;` must be decoded before any
        # text containing literal `&` could be misinterpreted, but since
        # we're only doing a flat replace it's safe in either order here.
        text = text.replace("&nbsp;", " ")
        text = text.replace("&amp;", "&")
        text = text.replace("&lt;", "<")
        text = text.replace("&gt;", ">")

    # Split into lines so we can scan top-down for the first reply marker.
    # We work line-by-line because reply chains are line-anchored — a
    # multi-line regex over the whole body would be much harder to reason
    # about.
    lines = text.split("\n")

    # Inner helper. Defined inside parse_email_body so it can close over
    # `lines` (it needs to look ahead by one for the From:/Sent: pair).
    def should_cut_from_line_index(i: int) -> bool:
        """Return True if email parsing should discard this line and everything after it.

        One-sentence summary: detects standard reply/signature dividers so only the newest message body is kept.

        Why it exists: Without cutting, the same Jira discussion or retro decision is re-ingested after every
        "Reply" — vector DB fills with duplicates and ranking surfaces stale paragraphs instead of the final answer.

        Decisions (first True wins in the outer scan; patterns here are mutually tested per line):
            1. ``^On .+ wrote:$`` — Gmail/Apple Mail style reply header; anchored so mid-line prose does not match.
            2. ``--- Original Message ---`` — Outlook delimiter; case-insensitive for localized templates.
            3. Lines starting with ``>`` — plaintext quotes; cutting once removes the entire remainder cheaply.
            4. ``^_{5,}`` — Outlook signature separator; five underscores avoids matching a single underscore in text.
            5. ``From:`` line only if next line starts with ``Sent:`` — avoids false cuts on legitimate "From:" sentences.

        Returns ``bool``. Outer loop sets ``cut_idx`` on first True. A mistaken True truncates the entire email to
        the first line — catastrophic for ingestion; a mistaken False keeps full threads — noisy RAG answers.

        Raises:
            Never raises.
        """

        # Pick the line we are currently inspecting. Stored in a local so
        # the regex calls below stay readable.
        line = lines[i]
        # "On Tuesday, ... wrote:" — the universal Gmail/Outlook reply
        # header. Anchored to start so it doesn't match if it's quoted
        # inline mid-sentence. `\s*$` tolerates trailing whitespace.
        if re.match(r"^On .+ wrote:\s*$", line):
            # RETURN: triggered when this line is the Gmail/Outlook reply
            # header. The outer scan uses this True as the cut index, so
            # everything from this line onwards (the quoted previous
            # email) is dropped before chunking. Without this cut we'd
            # ingest the same conversation N times, ballooning the index.
            return True
        # Older Outlook style. Case-insensitive match because the casing
        # of "Original Message" varies across locales.
        if re.match(r"^---\s*Original Message\s*---\s*$", line, re.I):
            # RETURN: triggered when an old-school Outlook divider
            # ("--- Original Message ---") is seen. Same downstream
            # effect as above: the outer scan cuts the body here so the
            # historical thread isn't re-ingested.
            return True
        # Any line starting with `>` is a quoted reply line. We cut at the
        # first one rather than skipping individual lines because the
        # quoted block always continues to the end of the message.
        if line.startswith(">"):
            # RETURN: triggered when we hit a `>` quoted-reply line. The
            # outer scan cuts here so the quoted history is excluded.
            # Cutting (rather than skipping) is important — `>` blocks
            # are interleaved with empty lines, and skipping individually
            # would keep noise we don't want.
            return True
        # Long underscore separators (`_____...`) are Outlook's signature/
        # reply divider. Five-or-more covers both five-char and longer
        # variants without flagging a single underscore in body text.
        if re.match(r"^_{5,}", line):
            # RETURN: triggered when an Outlook signature/reply underscore
            # bar is seen. The outer scan cuts here so signatures and the
            # reply chain after them are excluded from chunking. Without
            # this rule, every email would carry footer boilerplate that
            # leaks into retrieval results.
            return True
        # The "From: X / Sent: Y" pair is Outlook's plain-text reply
        # header. We require BOTH lines (not just "From:" alone) because
        # legitimate body text can start with "From:" — e.g. "From: the
        # whole team, thanks!". Looking ahead by one line gives us a
        # nearly zero false-positive rate.
        if line.startswith("From:") and i + 1 < len(lines):
            if lines[i + 1].lstrip().startswith("Sent:"):
                # RETURN: triggered when we see the From:/Sent: pair that
                # opens an Outlook plain-text reply. Outer scan cuts at
                # this line so the prior message is excluded. Returning
                # True on a bare "From:" without the Sent: lookahead
                # would chop out legitimate body text starting with
                # "From: the whole team...".
                return True
        # RETURN: triggered when the line is ordinary body text (no reply
        # marker). Outer scan continues to the next line. If this False
        # accidentally became a True, the email body would be truncated
        # at the very first line and the actual message would never
        # reach the chunker.
        return False

    # Scan for the first cut point. We use an index-based loop because the
    # helper needs the index for its lookahead.
    # `cut_idx` stays None until a reply divider fires — meaning "keep the entire body".
    # Using Optional[int] avoids magic sentinel (-1) that could be confused with a valid line index.
    cut_idx: int | None = None
    for i in range(len(lines)):
        if should_cut_from_line_index(i):
            cut_idx = i
            break

    # Truncate the body at the first reply marker (everything from that
    # line on is the previous message). If no marker was found, we keep
    # the entire body.
    if cut_idx is not None:
        lines = lines[:cut_idx]

    # Reassemble the trimmed body.
    text = "\n".join(lines)
    # Collapse runs of 3+ newlines to exactly 2. Email signatures and
    # mailing-list footers love to leave 5+ blank lines after the body,
    # which would create empty chunks downstream.
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Collapse runs of 2+ spaces to a single space. HTML-derived plain
    # text is full of double spaces from `&nbsp;` removal and indentation.
    text = re.sub(r" {2,}", " ", text)

    # Final strip removes leading/trailing whitespace that the collapse
    # rules above might leave behind — important so chunk content doesn't
    # start with a blank line.
    # RETURN: triggered after every cleanup step (HTML stripping, reply
    # truncation, whitespace normalisation) has run. The string returned
    # is the canonical "what the human actually wrote" version of the
    # email. Pipeline.py passes this string to chunk_email along with
    # sender/subject/date metadata; chunk_email decides whether the body
    # fits in one chunk or needs paragraph-level splitting. If we
    # accidentally returned the raw HTML or the un-trimmed thread, every
    # email chunk would be polluted with quoted history and embedded
    # garbage that hurts retrieval quality.
    return text.strip()

